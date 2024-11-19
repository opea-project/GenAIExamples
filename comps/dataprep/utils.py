# Copyright (C) 2024 Intel Corporation
# SPDX-License-Identifier: Apache-2.0

import base64
import errno
import functools
import json
import multiprocessing
import os
import re
import shutil
import signal
import subprocess
import tempfile
import timeit
import unicodedata
import urllib.parse
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Dict, List, Union
from urllib.parse import urlparse, urlunparse

import cairosvg
import cv2
import docx
import docx2txt
import fitz
import numpy as np
import pandas as pd
import pptx
import pytesseract
import requests
import yaml
from bs4 import BeautifulSoup
from langchain import LLMChain, PromptTemplate
from langchain_community.document_loaders import (
    UnstructuredHTMLLoader,
    UnstructuredImageLoader,
    UnstructuredMarkdownLoader,
    UnstructuredXMLLoader,
)
from langchain_community.llms import HuggingFaceEndpoint

from comps import CustomLogger

logger = CustomLogger("prepare_doc_util")
logflag = os.getenv("LOGFLAG", False)


class TimeoutError(Exception):
    pass


def timeout(seconds=10, error_message=os.strerror(errno.ETIME)):
    def decorator(func):
        def _handle_timeout(signum, frame):
            raise TimeoutError(error_message)

        @functools.wraps(func)
        def wrapper(*args, **kwargs):
            signal.signal(signal.SIGALRM, _handle_timeout)
            signal.alarm(seconds)
            try:
                result = func(*args, **kwargs)
            finally:
                signal.alarm(0)
            return result

        return wrapper

    return decorator


class Timer:
    level = 0
    viewer = None

    def __init__(self, name):
        self.name = name
        if Timer.viewer:
            Timer.viewer.display(f"{name} started ...")
        else:
            print(f"{name} started ...")

    def __enter__(self):
        self.start = timeit.default_timer()
        Timer.level += 1

    def __exit__(self, *a, **kw):
        Timer.level -= 1
        if Timer.viewer:
            Timer.viewer.display(f'{"  " * Timer.level}{self.name} took {timeit.default_timer() - self.start} sec')
        else:
            print(f'{"  " * Timer.level}{self.name} took {timeit.default_timer() - self.start} sec')


def get_separators():
    separators = [
        "\n\n",
        "\n",
        " ",
        ".",
        ",",
        "\u200b",  # Zero-width space
        "\uff0c",  # Fullwidth comma
        "\u3001",  # Ideographic comma
        "\uff0e",  # Fullwidth full stop
        "\u3002",  # Ideographic full stop
        "",
    ]
    return separators


def process_page(doc, idx):
    page = doc.load_page(idx)
    pagetext = page.get_text().strip()
    result = pagetext if pagetext.endswith(("!", "?", ".")) else pagetext + "."

    page_images = doc.get_page_images(idx)
    if page_images:
        for img_index, img in enumerate(page_images):
            xref = img[0]
            img_data = doc.extract_image(xref)
            img_bytes = img_data["image"]

            # process images
            img_array = cv2.imdecode(np.frombuffer(img_bytes, np.uint8), cv2.IMREAD_COLOR)
            img_result = pytesseract.image_to_string(img_array, lang="eng", config="--psm 6")

            # add results
            pageimg = img_result.strip()
            pageimg += "" if pageimg.endswith(("!", "?", ".")) else "."
            result += pageimg
    return result


def load_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    results = []

    with ThreadPoolExecutor(max_workers=4) as executor:
        futures = [executor.submit(process_page, doc, i) for i in range(doc.page_count)]
        for future in as_completed(futures):
            results.append(future.result())

    combined_result = "".join(results)
    return combined_result


def load_html(html_path):
    """Load the html file."""
    data_html = UnstructuredHTMLLoader(html_path).load()
    content = ""
    for ins in data_html:
        content += ins.page_content
    return content


def load_txt(txt_path):
    """Load txt file."""
    with open(txt_path, "r") as file:
        text = file.read()
    return text


def load_doc(doc_path):
    """Load doc file."""
    print("Converting doc file to docx file...")
    docx_path = doc_path + "x"
    subprocess.run(
        [
            "libreoffice",
            "--headless",
            "--invisible",
            "--convert-to",
            "docx",
            "--outdir",
            os.path.dirname(docx_path),
            doc_path,
        ],
        check=True,
    )
    print("Converted doc file to docx file.")
    text = load_docx(docx_path)
    os.remove(docx_path)
    return text


def load_docx(docx_path):
    """Load docx file."""
    doc = docx.Document(docx_path)
    text = ""
    # Save all 'rId:filenames' relationships in an dictionary and save the images if any.
    rid2img = {}
    for r in doc.part.rels.values():
        if isinstance(r._target, docx.parts.image.ImagePart):
            rid2img[r.rId] = os.path.basename(r._target.partname)
    if rid2img:
        save_path = tempfile.mkdtemp()
        docx2txt.process(docx_path, save_path)
    for paragraph in doc.paragraphs:
        if hasattr(paragraph, "text"):
            text += paragraph.text + "\n"
        if "graphicData" in paragraph._p.xml:
            for rid in rid2img:
                if rid in paragraph._p.xml:
                    img_path = os.path.join(save_path, rid2img[rid])
                    img_text = load_image(img_path)
                    if img_text:
                        text += img_text + "\n"
    if rid2img:
        shutil.rmtree(save_path)
    return text


def load_ppt(ppt_path):
    """Load ppt file."""
    print("Converting ppt file to pptx file...")
    pptx_path = ppt_path + "x"
    subprocess.run(
        [
            "libreoffice",
            "--headless",
            "--invisible",
            "--convert-to",
            "docx",
            "--outdir",
            os.path.dirname(pptx_path),
            ppt_path,
        ],
        check=True,
    )
    print("Converted ppt file to pptx file.")
    text = load_pptx(pptx_path)
    os.remove(pptx_path)
    return text


def load_pptx(pptx_path):
    """Load pptx file."""
    text = ""
    prs = pptx.Presentation(pptx_path)
    for slide in prs.slides:
        for shape in sorted(slide.shapes, key=lambda shape: (shape.top, shape.left)):
            if shape.has_text_frame:
                if shape.text:
                    text += shape.text + "\n"
            if shape.has_table:
                table_contents = "\n".join(
                    [
                        "\t".join([(cell.text if hasattr(cell, "text") else "") for cell in row.cells])
                        for row in shape.table.rows
                        if hasattr(row, "cells")
                    ]
                )
                if table_contents:
                    text += table_contents + "\n"
            if hasattr(shape, "image") and hasattr(shape.image, "blob"):
                img_path = f"./{shape.image.filename}"
                with open(img_path, "wb") as f:
                    f.write(shape.image.blob)
                img_text = load_image(img_path)
                if img_text:
                    text += img_text + "\n"
                os.remove(img_path)
    return text


def load_md(md_path):
    """Load md file."""
    loader = UnstructuredMarkdownLoader(md_path)
    text = loader.load()[0].page_content
    return text


def load_xml(xml_path):
    """Load xml file."""
    loader = UnstructuredXMLLoader(xml_path)
    text = loader.load()[0].page_content
    return text


def load_json(json_path):
    """Load and process json file."""
    with open(json_path, "r") as file:
        data = json.load(file)
    content_list = [json.dumps(item) for item in data]
    return content_list


def load_jsonl(jsonl_path):
    """Load and process jsonl file."""
    content_list = []
    with open(jsonl_path, "r") as file:
        for line in file:
            json_obj = json.loads(line)
            content_list.append(json_obj)
    return content_list


def load_yaml(yaml_path):
    """Load and process yaml file."""
    with open(yaml_path, "r") as file:
        data = yaml.safe_load(file)
    return yaml.dump(data)


def load_xlsx(input_path):
    """Load and process xlsx file."""
    df = pd.read_excel(input_path)
    content_list = df.apply(lambda row: ", ".join(row.astype(str)), axis=1).tolist()
    return content_list


def load_csv(input_path):
    """Load the csv file."""
    df = pd.read_csv(input_path)
    content_list = df.apply(lambda row: ", ".join(row.astype(str)), axis=1).tolist()
    return content_list


def load_image(image_path):
    """Load the image file."""
    if os.getenv("SUMMARIZE_IMAGE_VIA_LVM", None) == "1":
        query = "Please summarize this image."
        image_b64_str = base64.b64encode(open(image_path, "rb").read()).decode()
        response = requests.post(
            "http://localhost:9399/v1/lvm",
            data=json.dumps({"image": image_b64_str, "prompt": query}),
            headers={"Content-Type": "application/json"},
            proxies={"http": None},
        )
        return response.json()["text"].strip()
    loader = UnstructuredImageLoader(image_path)
    text = loader.load()[0].page_content
    return text.strip()


def load_svg(svg_path):
    """Load the svg file."""
    png_path = svg_path.replace(".svg", ".png")
    cairosvg.svg2png(url=svg_path, write_to=png_path)
    text = load_image(png_path)
    os.remove(png_path)
    return text


def document_loader(doc_path):
    if doc_path.endswith(".pdf"):
        return load_pdf(doc_path)
    elif doc_path.endswith(".html"):
        return load_html(doc_path)
    elif doc_path.endswith(".txt"):
        return load_txt(doc_path)
    elif doc_path.endswith(".doc"):
        return load_doc(doc_path)
    elif doc_path.endswith(".docx"):
        return load_docx(doc_path)
    elif doc_path.endswith(".ppt"):
        return load_ppt(doc_path)
    elif doc_path.endswith(".pptx"):
        return load_pptx(doc_path)
    elif doc_path.endswith(".md"):
        return load_md(doc_path)
    elif doc_path.endswith(".xml"):
        return load_xml(doc_path)
    elif doc_path.endswith(".json"):
        return load_json(doc_path)
    elif doc_path.endswith(".jsonl"):
        return load_jsonl(doc_path)
    elif doc_path.endswith(".yaml"):
        return load_yaml(doc_path)
    elif doc_path.endswith(".xlsx") or doc_path.endswith(".xls"):
        return load_xlsx(doc_path)
    elif doc_path.endswith(".csv"):
        return load_csv(doc_path)
    elif (
        doc_path.endswith(".tiff")
        or doc_path.endswith(".jpg")
        or doc_path.endswith(".jpeg")
        or doc_path.endswith(".png")
    ):
        return load_image(doc_path)
    elif doc_path.endswith(".svg"):
        return load_svg(doc_path)
    else:
        raise NotImplementedError(
            "Current only support pdf, html, txt, doc, docx, pptx, ppt, md, xml"
            + ", json, jsonl, yaml, xlsx, xls, csv, tiff and svg format."
        )


class Crawler:
    def __init__(self, pool=None):
        if pool:
            assert isinstance(pool, (str, list, tuple)), "url pool should be str, list or tuple"
        self.pool = pool
        self.headers = {
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng, \
            */*;q=0.8,application/signed-exchange;v=b3;q=0.7",
            "Accept-Encoding": "gzip, deflate, br",
            "Accept-Language": "en-US,en;q=0.9,zh-CN;q=0.8,zh;q=0.7",
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, \
            like Gecko) Chrome/113.0.0.0 Safari/537.36",
        }
        self.fetched_pool = set()

    def get_sublinks(self, soup):
        sublinks = []
        for links in soup.find_all("a"):
            sublinks.append(str(links.get("href")))
        return sublinks

    def get_hyperlink(self, soup, base_url):
        sublinks = []
        for links in soup.find_all("a"):
            link = str(links.get("href"))
            if link.startswith("#") or link is None or link == "None":
                continue
            suffix = link.split("/")[-1]
            if "." in suffix and suffix.split(".")[-1] not in ["html", "htmld"]:
                continue
            link_parse = urlparse(link)
            base_url_parse = urlparse(base_url)
            if link_parse.path == "":
                continue
            if link_parse.netloc != "":
                # keep crawler works in the same domain
                if link_parse.netloc != base_url_parse.netloc:
                    continue
                sublinks.append(link)
            else:
                sublinks.append(
                    urlunparse(
                        (
                            base_url_parse.scheme,
                            base_url_parse.netloc,
                            link_parse.path,
                            link_parse.params,
                            link_parse.query,
                            link_parse.fragment,
                        )
                    )
                )
        return sublinks

    def fetch(self, url, headers=None, max_times=5):
        if not headers:
            headers = self.headers
        while max_times:
            parsed_url = urlparse(url)
            if not parsed_url.scheme:
                url = "http://" + url
            if logflag:
                logger.info("start fetch %s..." % url)
            try:
                response = requests.get(url, headers=headers, verify=True)
                if response.status_code != 200:
                    print("fail to fetch %s, response status code: %s", url, response.status_code)
                else:
                    # Extract charset from the Content-Type header
                    content_type = response.headers.get("Content-Type", "").lower()
                    if "charset=" in content_type:
                        # Extract charset value from the content-type header
                        charset = content_type.split("charset=")[-1].strip()
                        response.encoding = charset
                        if logflag:
                            logger.info(f"Charset detected and set: {response.encoding}")
                    else:
                        import re

                        # Extract charset from the response HTML content
                        charset_from_meta = None
                        # Check for <meta charset="...">
                        match = re.search(r'<meta\s+charset=["\']?([^"\'>]+)["\']?', response.text, re.IGNORECASE)
                        if match:
                            charset_from_meta = match.group(1)
                        # Check for <meta http-equiv="Content-Type" content="...; charset=...">
                        if not charset_from_meta:
                            match = re.search(
                                r'<meta\s+http-equiv=["\']?content-type["\']?\s+content=["\']?[^"\']*charset=([^"\'>]+)["\']?',
                                response.text,
                                re.IGNORECASE,
                            )
                            if match:
                                charset_from_meta = match.group(1)
                        if charset_from_meta:
                            response.encoding = charset_from_meta
                            if logflag:
                                logger.info(f"Charset detected and set from meta tag: {response.encoding}")
                        else:
                            # Fallback to default encoding
                            response.encoding = "utf-8"
                            if logflag:
                                logger.info("Charset not specified, using default utf-8")
                    return response
            except Exception as e:
                print("fail to fetch %s, caused by %s", url, e)
                raise Exception(e)
            max_times -= 1
        return None

    def process_work(self, sub_url, work):
        response = self.fetch(sub_url)
        if response is None:
            return []
        self.fetched_pool.add(sub_url)
        soup = self.parse(response.text)
        base_url = self.get_base_url(sub_url)
        sublinks = self.get_hyperlink(soup, base_url)
        if work:
            work(sub_url, soup)
        return sublinks

    def crawl(self, pool, work=None, max_depth=10, workers=10):
        url_pool = set()
        for url in pool:
            base_url = self.get_base_url(url)
            response = self.fetch(url)
            soup = self.parse(response.text)
            sublinks = self.get_hyperlink(soup, base_url)
            self.fetched_pool.add(url)
            url_pool.update(sublinks)
            depth = 0
            while len(url_pool) > 0 and depth < max_depth:
                print("current depth %s...", depth)
                mp = multiprocessing.Pool(processes=workers)
                results = []
                for sub_url in url_pool:
                    if sub_url not in self.fetched_pool:
                        results.append(mp.apply_async(self.process_work, (sub_url, work)))
                mp.close()
                mp.join()
                url_pool = set()
                for result in results:
                    sublinks = result.get()
                    url_pool.update(sublinks)
                depth += 1

    def parse(self, html_doc):
        soup = BeautifulSoup(html_doc, "lxml")
        return soup

    def download(self, url, file_name):
        print("download %s into %s...", url, file_name)
        try:
            r = requests.get(url, stream=True, headers=self.headers, verify=True)
            f = open(file_name, "wb")
            for chunk in r.iter_content(chunk_size=512):
                if chunk:
                    f.write(chunk)
        except Exception as e:
            print("fail to download %s, caused by %s", url, e)

    def get_base_url(self, url):
        result = urlparse(url)
        return urlunparse((result.scheme, result.netloc, "", "", "", ""))

    def clean_text(self, text):
        text = text.strip().replace("\r", "\n")
        text = re.sub(" +", " ", text)
        text = re.sub("\n+", "\n", text)
        text = text.split("\n")
        return "\n".join([i for i in text if i and i != " "])


def uni_pro(text):
    """Check if the character is ASCII or falls in the category of non-spacing marks."""
    normalized_text = unicodedata.normalize("NFKD", text)
    filtered_text = ""
    for char in normalized_text:
        if ord(char) < 128 or unicodedata.category(char) == "Mn":
            filtered_text += char
    return filtered_text


def load_html_data(url):
    crawler = Crawler()
    res = crawler.fetch(url)
    if res is None:
        return None
    soup = crawler.parse(res.text)
    all_text = crawler.clean_text(soup.select_one("body").text)
    main_content = ""
    for element_name in ["main", "container"]:
        main_block = None
        if soup.select(f".{element_name}"):
            main_block = soup.select(f".{element_name}")
        elif soup.select(f"#{element_name}"):
            main_block = soup.select(f"#{element_name}")
        if main_block:
            for element in main_block:
                text = crawler.clean_text(element.text)
                if text not in main_content:
                    main_content += f"\n{text}"
            main_content = crawler.clean_text(main_content)
    main_content = all_text if main_content == "" else main_content
    main_content = main_content.replace("\n", "")
    main_content = main_content.replace("\n\n", "")
    main_content = re.sub(r"\s+", " ", main_content)
    if logflag:
        logger.info("main_content=[%s]" % main_content)

    return main_content


def parse_html(input):
    """Parse the uploaded file."""
    chucks = []
    for link in input:
        if re.match(r"^https?:/{2}\w.+$", link):
            content = load_html_data(link)
            if content is None:
                continue
            chuck = [[content.strip(), link]]
            chucks += chuck
        else:
            print("The given link/str {} cannot be parsed.".format(link))

    return chucks


def load_html_content(links, chunk_size=1500, chunk_overlap=50):
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    from langchain_community.document_loaders import AsyncHtmlLoader
    from langchain_community.document_transformers import Html2TextTransformer

    loader = AsyncHtmlLoader(links, ignore_load_errors=True, trust_env=True)
    docs = loader.load()
    html2text = Html2TextTransformer()
    docs = list(html2text.transform_documents(docs))
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    docs = text_splitter.split_documents(docs)
    return docs


def parse_html_new(input, chunk_size, chunk_overlap):
    docs = load_html_content(input, chunk_size, chunk_overlap)
    html_content = ""
    for doc in docs:
        html_content += doc.page_content + "\n"
    return html_content


def get_tables_result(pdf_path, table_strategy):
    """Extract tables information from pdf file."""
    if table_strategy == "fast":
        return None

    from unstructured.documents.elements import FigureCaption
    from unstructured.partition.pdf import partition_pdf

    tables_result = []
    raw_pdf_elements = partition_pdf(
        filename=pdf_path,
        infer_table_structure=True,
    )
    tables = [el for el in raw_pdf_elements if el.category == "Table"]
    for table in tables:
        table_coords = table.metadata.coordinates.points
        content = table.metadata.text_as_html
        table_page_number = table.metadata.page_number
        min_distance = float("inf")
        table_summary = None
        if table_strategy == "hq":
            for element in raw_pdf_elements:
                if isinstance(element, FigureCaption) or element.text.startswith("Tab"):
                    caption_page_number = element.metadata.page_number
                    caption_coords = element.metadata.coordinates.points
                    related, y_distance = get_relation(
                        table_coords, caption_coords, table_page_number, caption_page_number
                    )
                    if related:
                        if y_distance < min_distance:
                            min_distance = y_distance
                            table_summary = element.text
            if table_summary is None:
                parent_id = table.metadata.parent_id
                for element in raw_pdf_elements:
                    if element.id == parent_id:
                        table_summary = element.text
                        break
        elif table_strategy == "llm":
            table_summary = llm_generate(content)
            table_summary = table_summary.lstrip("\n ")
        elif table_strategy is None:
            table_summary = None
        if table_summary is None:
            text = f"[Table: {content}]"
        else:
            text = f"|Table: [Summary: {table_summary}], [Content: {content}]|"
        tables_result.append(text)
    return tables_result


def llm_generate(content):
    llm_endpoint = os.getenv("TGI_LLM_ENDPOINT", "http://localhost:8080")
    llm = HuggingFaceEndpoint(
        endpoint_url=llm_endpoint,
        max_new_tokens=1000,
        top_k=40,
        top_p=0.9,
        temperature=0.8,
        streaming=False,
        num_beams=2,
        num_return_sequences=2,
        use_cache=True,
        timeout=600,
    )

    table_summary_template = """
    Task: Your task is to give a concise summary of the table. \
    The summary should cover the overall table structure and all detailed information of the table. \
    The table will be given in html format. Summarize the table below.
    ---
    ### Table:
    {table_content}
    ---
    ### Generated Summary:
    """

    prompt = PromptTemplate(template=table_summary_template, input_variables=["table_content"])

    llm_chain = LLMChain(prompt=prompt, llm=llm)

    response = llm_chain.invoke(content)
    response = response["text"]
    print("response", response)
    return response


def get_relation(table_coords, caption_coords, table_page_number, caption_page_number, threshold=100):
    """Get the relation of a pair of table and caption."""
    same_page = table_page_number == caption_page_number
    x_overlap = (min(table_coords[2][0], caption_coords[2][0]) - max(table_coords[0][0], caption_coords[0][0])) > 0
    if table_coords[0][1] - caption_coords[1][1] >= 0:
        y_distance = table_coords[0][1] - caption_coords[1][1]
    elif caption_coords[0][1] - table_coords[1][1] >= 0:
        y_distance = caption_coords[0][1] - table_coords[1][1]
    else:
        y_distance = 0
    y_close = y_distance < threshold
    return same_page and x_overlap and y_close, y_distance


def create_upload_folder(upload_path):
    if not os.path.exists(upload_path):
        Path(upload_path).mkdir(parents=True, exist_ok=True)


def encode_filename(filename):
    return urllib.parse.quote(filename, safe="")


def decode_filename(encoded_filename):
    return urllib.parse.unquote(encoded_filename)


async def save_content_to_local_disk(save_path: str, content):
    save_path = Path(save_path)
    try:
        if isinstance(content, str):
            with open(save_path, "w", encoding="utf-8") as file:
                file.write(content)
        else:
            with save_path.open("wb") as fout:
                content = await content.read()
                fout.write(content)
    except Exception as e:
        print(f"Write file failed. Exception: {e}")
        raise Exception(status_code=500, detail=f"Write file {save_path} failed. Exception: {e}")


def get_file_structure(root_path: str, parent_path: str = "") -> List[Dict[str, Union[str, List]]]:
    result = []
    for path in os.listdir(root_path):
        complete_path = parent_path + "/" + path if parent_path else path
        file_path = root_path + "/" + path
        p = Path(file_path)
        # append file into result
        if p.is_file():
            file_dict = {
                "name": decode_filename(path),
                "id": decode_filename(complete_path),
                "type": "File",
                "parent": "",
            }
            result.append(file_dict)
        else:
            # append folder and inner files/folders into result using recursive function
            folder_dict = {
                "name": decode_filename(path),
                "id": decode_filename(complete_path),
                "type": "Directory",
                "children": get_file_structure(file_path, complete_path),
                "parent": "",
            }
            result.append(folder_dict)

    return result


def format_search_results(response, file_list: list):
    for i in range(1, len(response), 2):
        file_name = response[i].decode()[5:]
        file_dict = {
            "name": decode_filename(file_name),
            "id": decode_filename(file_name),
            "type": "File",
            "parent": "",
        }
        file_list.append(file_dict)
    return file_list


def remove_folder_with_ignore(folder_path: str, except_patterns: List = []):
    """Remove the specific folder, and ignore some files/folders.

    :param folder_path: file path to delete
    :param except_patterns: files/folder name to ignore
    """
    print(f"except patterns: {except_patterns}")
    for root, dirs, files in os.walk(folder_path, topdown=False):
        for name in files:
            # delete files except ones that match patterns
            file_path = os.path.join(root, name)
            if except_patterns != [] and any(pattern in file_path for pattern in except_patterns):
                continue
            os.remove(file_path)

        # delete empty folder
        for name in dirs:
            dir_path = os.path.join(root, name)
            # delete folders except ones that match patterns
            if except_patterns != [] and any(pattern in dir_path for pattern in except_patterns):
                continue
            if not os.listdir(dir_path):
                os.rmdir(dir_path)
