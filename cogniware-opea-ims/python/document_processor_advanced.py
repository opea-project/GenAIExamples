"""
Cogniware Core - Advanced Document Processor
Support for PDF, DOCX, XLS, XLSX, PPT, PPTX and file uploads
"""

import os
import io
import base64
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, List

# Try to import document processing libraries
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

class AdvancedDocumentProcessor:
    """Process multiple document types including uploads"""
    
    def __init__(self, documents_dir):
        self.documents_dir = Path(documents_dir)
        self.documents_dir.mkdir(exist_ok=True)
        self.supported_types = {
            'pdf': PDF_AVAILABLE,
            'docx': DOCX_AVAILABLE,
            'doc': DOCX_AVAILABLE,
            'xlsx': EXCEL_AVAILABLE,
            'xls': EXCEL_AVAILABLE,
            'pptx': PPTX_AVAILABLE,
            'ppt': PPTX_AVAILABLE,
            'txt': True,
            'md': True,
            'csv': True,
            'json': True
        }
    
    def get_supported_formats(self) -> dict:
        """Get list of supported formats"""
        return {
            'success': True,
            'formats': {
                'documents': {
                    'pdf': PDF_AVAILABLE,
                    'docx': DOCX_AVAILABLE,
                    'doc': DOCX_AVAILABLE
                },
                'spreadsheets': {
                    'xlsx': EXCEL_AVAILABLE,
                    'xls': EXCEL_AVAILABLE,
                    'csv': True
                },
                'presentations': {
                    'pptx': PPTX_AVAILABLE,
                    'ppt': PPTX_AVAILABLE
                },
                'text': {
                    'txt': True,
                    'md': True,
                    'json': True
                }
            },
            'install_command': {
                'pdf': 'pip install PyPDF2',
                'docx': 'pip install python-docx',
                'excel': 'pip install openpyxl',
                'powerpoint': 'pip install python-pptx'
            }
        }
    
    def upload_document(self, filename: str, file_data: bytes, doc_type: str = None) -> dict:
        """Upload and save document"""
        try:
            # Determine type from filename if not provided
            if not doc_type:
                doc_type = filename.split('.')[-1].lower()
            
            # Check if format is supported
            if doc_type not in self.supported_types:
                return {
                    'success': False,
                    'error': f'Unsupported file type: {doc_type}'
                }
            
            # Save file
            filepath = self.documents_dir / filename
            filepath.write_bytes(file_data)
            
            return {
                'success': True,
                'filename': filename,
                'filepath': str(filepath),
                'size': len(file_data),
                'type': doc_type
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}
    
    def upload_document_base64(self, filename: str, base64_data: str, doc_type: str = None) -> dict:
        """Upload document from base64 encoded data"""
        try:
            # Decode base64
            file_data = base64.b64decode(base64_data)
            return self.upload_document(filename, file_data, doc_type)
        except Exception as e:
            return {'success': False, 'error': str(e)}
    
    def extract_text_from_pdf(self, filepath: Path) -> dict:
        """Extract text from PDF"""
        if not PDF_AVAILABLE:
            return {'success': False, 'error': 'PyPDF2 not installed. Run: pip install PyPDF2'}
        
        try:
            text_content = []
            with open(filepath, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                page_count = len(pdf_reader.pages)
                
                for page_num in range(page_count):
                    page = pdf_reader.pages[page_num]
                    text_content.append(page.extract_text())
            
            full_text = '\n'.join(text_content)
            
            return {
                'success': True,
                'text': full_text,
                'page_count': page_count,
                'word_count': len(full_text.split()),
                'char_count': len(full_text)
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}
    
    def extract_text_from_docx(self, filepath: Path) -> dict:
        """Extract text from DOCX"""
        if not DOCX_AVAILABLE:
            return {'success': False, 'error': 'python-docx not installed. Run: pip install python-docx'}
        
        try:
            doc = DocxDocument(filepath)
            paragraphs = [p.text for p in doc.paragraphs]
            full_text = '\n'.join(paragraphs)
            
            return {
                'success': True,
                'text': full_text,
                'paragraph_count': len(paragraphs),
                'word_count': len(full_text.split()),
                'char_count': len(full_text)
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}
    
    def extract_data_from_excel(self, filepath: Path) -> dict:
        """Extract data from Excel"""
        if not EXCEL_AVAILABLE:
            return {'success': False, 'error': 'openpyxl not installed. Run: pip install openpyxl'}
        
        try:
            workbook = openpyxl.load_workbook(filepath)
            sheets_data = {}
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                data = []
                
                for row in sheet.iter_rows(values_only=True):
                    data.append(list(row))
                
                sheets_data[sheet_name] = {
                    'data': data,
                    'row_count': len(data),
                    'column_count': len(data[0]) if data else 0
                }
            
            return {
                'success': True,
                'sheets': sheets_data,
                'sheet_count': len(sheets_data)
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}
    
    def extract_text_from_pptx(self, filepath: Path) -> dict:
        """Extract text from PowerPoint"""
        if not PPTX_AVAILABLE:
            return {'success': False, 'error': 'python-pptx not installed. Run: pip install python-pptx'}
        
        try:
            prs = Presentation(filepath)
            slides_text = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        slide_text.append(shape.text)
                
                slides_text.append({
                    'slide_number': slide_num + 1,
                    'text': '\n'.join(slide_text)
                })
            
            full_text = '\n\n'.join([s['text'] for s in slides_text])
            
            return {
                'success': True,
                'text': full_text,
                'slides': slides_text,
                'slide_count': len(slides_text),
                'word_count': len(full_text.split()),
                'char_count': len(full_text)
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}
    
    def process_document(self, filename: str) -> dict:
        """Process any supported document type"""
        try:
            # Find document
            docs = list(self.documents_dir.glob(f"{filename}*"))
            if not docs:
                # Try exact filename
                filepath = self.documents_dir / filename
                if not filepath.exists():
                    return {'success': False, 'error': 'Document not found'}
            else:
                filepath = docs[0]
            
            # Get file extension
            file_ext = filepath.suffix.lower().replace('.', '')
            
            # Extract based on type
            result = {'success': True, 'filename': filepath.name, 'type': file_ext}
            
            if file_ext == 'pdf':
                extraction = self.extract_text_from_pdf(filepath)
                result.update(extraction)
            
            elif file_ext in ['docx', 'doc']:
                extraction = self.extract_text_from_docx(filepath)
                result.update(extraction)
            
            elif file_ext in ['xlsx', 'xls']:
                extraction = self.extract_data_from_excel(filepath)
                result.update(extraction)
            
            elif file_ext in ['pptx', 'ppt']:
                extraction = self.extract_text_from_pptx(filepath)
                result.update(extraction)
            
            elif file_ext in ['txt', 'md', 'csv', 'json']:
                # Plain text processing
                content = filepath.read_text()
                result['text'] = content
                result['line_count'] = len(content.split('\n'))
                result['word_count'] = len(content.split())
                result['char_count'] = len(content)
            
            else:
                return {'success': False, 'error': f'Unsupported file type: {file_ext}'}
            
            # Add file metadata
            stat = filepath.stat()
            result['size_bytes'] = stat.st_size
            result['modified'] = datetime.fromtimestamp(stat.st_mtime).isoformat()
            result['filepath'] = str(filepath)
            
            return result
            
        except Exception as e:
            return {'success': False, 'error': str(e)}
    
    def search_in_document(self, filename: str, query: str) -> dict:
        """Search within a specific document"""
        try:
            # Process document first
            result = self.process_document(filename)
            
            if not result.get('success'):
                return result
            
            # Get text content
            text = result.get('text', '')
            
            if not text:
                return {'success': False, 'error': 'No text content found in document'}
            
            # Search for query
            matches = []
            lines = text.split('\n')
            
            for line_num, line in enumerate(lines):
                if query.lower() in line.lower():
                    matches.append({
                        'line_number': line_num + 1,
                        'line_text': line.strip(),
                        'context': lines[max(0, line_num-1):min(len(lines), line_num+2)]
                    })
            
            return {
                'success': True,
                'document': filename,
                'query': query,
                'match_count': len(matches),
                'matches': matches[:10]  # Limit to first 10 matches
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

# Create global instance
advanced_doc_processor = None

def get_advanced_processor(documents_dir):
    global advanced_doc_processor
    if advanced_doc_processor is None:
        advanced_doc_processor = AdvancedDocumentProcessor(documents_dir)
    return advanced_doc_processor

